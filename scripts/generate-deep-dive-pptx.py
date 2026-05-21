#!/usr/bin/env python3
"""Generate an editable PPTX Deep-Dive Deck for ML Container Creator."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# AWS brand colors
DARK = RGBColor(0x23, 0x2F, 0x3E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, left, top, width, height, text, size=18, color=DARK,
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return tf


def add_bullet_frame(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return tf


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'Calibri'
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                p.font.name = 'Calibri'
    return tbl


def orange_bar(slide, top=0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top),
        prs.slide_width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===========================================================================
# SLIDE 1: Title
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
orange_bar(s, 7.2)
text_box(s, 1, 1.5, 11, 1.2, 'ML Container Creator', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1, 2.8, 11, 0.8, 'Eliminating the BYOC Tax for LLM Deployments on SageMaker', size=24, color=ORANGE, align=PP_ALIGN.CENTER)
text_box(s, 1, 4.2, 11, 0.5, 'Internal Deep Dive \u00b7 May 2026', size=18, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)
text_box(s, 1, 5.2, 11, 0.5, 'npx @aws/ml-container-creator', size=16, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
text_box(s, 1, 6.0, 11, 0.5, 'Open Source \u00b7 Apache 2.0 \u00b7 github.com/awslabs/ml-container-creator', size=14, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
speaker_notes(s, 'Open with: "How many of you have helped a customer deploy an LLM on a SageMaker endpoint? How long did it take them to get their first successful inference?"\n\nSet the tone: this tool was born from real customer pain, battle-tested against real platform issues.\n\nMention: open-source, AWS Labs, Apache-2.0.')

# ===========================================================================
# SLIDE 2: The Customer Pain
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'The Customer Pain', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'Deploying an LLM on SageMaker requires getting 6 layers right', size=18, color=GRAY)
add_table(s, 0.8, 1.8, 11.7, 3.5,
    ['Layer', 'Decisions', 'Failure Mode'],
    [
        ['Container', 'Base image, CUDA, compat shim, serve script', 'Silent crash, OOM'],
        ['Model Server', 'vLLM version, CLI args, env vars, boolean flags', 'Invalid args, startup hang'],
        ['Infrastructure', 'Instance type, GPU count, tensor parallelism', 'Scheduling failure'],
        ['SageMaker', 'IC config, RoutingConfig, endpoint config', 'IC never goes InService'],
        ['Observability', 'Log groups, log streams, CW routing', 'Zero visibility into failures'],
        ['Adapters', 'LoRA paths, max rank, hot-swap lifecycle', 'Inference errors'],
    ],
    col_widths=[2.2, 5.5, 4.0]
)
text_box(s, 0.8, 5.8, 11, 0.5, 'Teams spend 2-5 days on container boilerplate before testing inference.', size=16, color=GRAY, bold=True)
speaker_notes(s, "This isn't hypothetical. I just spent a week deploying Gemma 4 31B. Seven distinct configuration issues, each one a silent failure.\n\nThe customer doesn't know about RoutingConfig. They don't know vLLM v0.20.2 deprecated env vars. They don't know NCCL hangs on PCIe topologies.\n\nTransition: Let me show you what those failures actually look like.")

# ===========================================================================
# SLIDE 3: The Fragility Landscape
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Real Issues Found Deploying One Model', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'Gemma 4 31B on SageMaker, May 2026 \u2014 7 distinct failures', size=18, color=GRAY)
add_table(s, 0.8, 1.8, 11.7, 3.8,
    ['#', 'Issue', 'Symptom'],
    [
        ['1', '--enable-lora true', 'vLLM crash: unrecognized argument'],
        ['2', 'VLLM_ENABLE_CUDA_COMPATIBILITY as CLI arg', 'Unknown argument error'],
        ['3', 'Missing RoutingConfig', 'IC stuck in Creating forever'],
        ['4', 'IC log routing', 'Log group created, zero log streams'],
        ['5', 'max_model_len=262144', 'OOM on 4\u00d724GB GPUs'],
        ['6', 'NCCL on g6 (L4, PCIe)', 'Hang before model load'],
        ['7', 'vLLM v0.20.2 env vars deprecated', 'Breaking change, no migration path'],
    ],
    col_widths=[0.6, 5.0, 6.1]
)
text_box(s, 0.8, 6.0, 11, 0.5, 'Every one invisible until runtime, 15 minutes into a deploy cycle.', size=16, color=ORANGE, bold=True)
speaker_notes(s, "Each of these cost 1-4 hours to diagnose. A customer without deep SageMaker + vLLM expertise would be stuck for days.\n\nThis is why MCC exists \u2014 to encode these lessons so the next customer doesn't hit them.\n\nTransition: So what does the solution look like?")

# ===========================================================================
# SLIDE 4: What MCC Does
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'What ML Container Creator Does', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.5, 'A Code Generator, Not a Runtime Framework', size=22, color=ORANGE, bold=True)
text_box(s, 0.8, 1.9, 5.0, 0.4, 'You provide:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 2.4, 5.0, 2.0, [
    '\u2022  Model name',
    '\u2022  Serving backend (vLLM, SGLang, Triton, ...)',
    '\u2022  Instance type',
    '\u2022  Deployment target',
], size=15, color=GRAY)
text_box(s, 7.0, 1.9, 5.5, 0.4, 'MCC generates:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 7.0, 2.4, 5.5, 2.0, [
    '\u2022  Dockerfile (multi-stage, optimized)',
    '\u2022  code/serve (entrypoint script)',
    '\u2022  do/ lifecycle scripts (20+)',
    '\u2022  Tests, IAM docs, README',
], size=15, color=GRAY)
text_box(s, 0.8, 4.8, 11, 0.4, 'Key principles:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 5.3, 11, 1.5, [
    '\u2022  You own the output \u2014 no runtime dependency on MCC',
    '\u2022  No agent in your container, no lock-in',
    '\u2022  Opinionated defaults, full escape hatches',
], size=15)
speaker_notes(s, 'Think create-react-app for ML containers.\n\nAfter generation, MCC is done. Your project has zero dependency on it. You can modify any file.\n\nAnalogy: "A senior ML engineer who\'s deployed 100 models, writing your boilerplate for you."')

# ===========================================================================
# SLIDE 5: Architecture
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Architecture', size=32, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 1.3, 11, 4.5, [
    'CLI \u2192 Config Manager (8-level precedence) \u2192 Template Engine (EJS)',
    '',
    '6 MCP Servers provide intelligent defaults:',
    '    instance-sizer, region-picker, base-image-picker,',
    '    model-picker, hyperpod-cluster-picker, endpoint-picker',
    '',
    '3 Catalogs validate every combination:',
    '    model-servers.json \u2014 14 server versions, base images, CUDA compat',
    '    models.json \u2014 Popular models, parameter counts, recommended instances',
    '    instances.json \u2014 48 types, GPU specs, memory, cost tiers',
    '',
    'Output: Complete project with Dockerfile + code/ + do/ + tests',
], size=15)
speaker_notes(s, 'The MCP servers are the intelligence layer. They know that Gemma 4 31B needs TP=4 on g5.12xlarge.\n\nThe catalogs prevent invalid combinations at generation time, not deploy time.\n\nThe 8-level precedence means: CLI flags > env vars > config file > MCP responses > defaults. This enables both interactive exploration and fully automated CI/CD.')

# ===========================================================================
# SLIDE 6: Supported Configurations
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Supported Configurations', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.0, 11, 0.4, '4 Architectures \u00d7 15 Configs \u00d7 4 Targets', size=18, color=GRAY)
add_table(s, 0.8, 1.6, 11.7, 2.5,
    ['Architecture', 'Backends', 'Use Case'],
    [
        ['HTTP (2)', 'Flask, FastAPI', 'Traditional ML (sklearn, XGBoost, TF)'],
        ['Transformers (5)', 'vLLM, SGLang, TRT-LLM, LMI, DJL', 'LLM serving'],
        ['Triton (7)', 'FIL, ONNX, TF, PyTorch, vLLM, TRT-LLM, Python', 'Multi-framework'],
        ['Diffusors (1)', 'vLLM-Omni', 'Image generation'],
    ],
    col_widths=[2.5, 5.5, 3.7]
)
add_table(s, 0.8, 4.6, 11.7, 2.2,
    ['Deployment Target', 'Key Feature'],
    [
        ['Managed Inference', 'Inference Components, multi-model, auto-scaling'],
        ['Async Inference', 'S3 output, SNS notifications, large payloads'],
        ['Batch Transform', 'S3-to-S3 bulk processing'],
        ['HyperPod EKS', 'Kubernetes, GPU scheduling, heterogeneous clusters'],
    ],
    col_widths=[3.0, 8.7]
)
speaker_notes(s, 'Same Dockerfile works across all 4 targets. The do/ scripts handle infrastructure differences.\n\n15 tested combinations means 15 Dockerfiles + serve scripts that we know work.\n\nAsk: "Which of these architectures are you currently using or evaluating?"')

# ===========================================================================
# SLIDE 7: The do/ Framework
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'The do/ Framework', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.0, 11, 0.4, '20+ Lifecycle Scripts \u2014 Your Entire Deployment Workflow', size=18, color=GRAY)
add_table(s, 0.8, 1.6, 11.7, 3.8,
    ['Category', 'Scripts', 'Purpose'],
    [
        ['Build', 'build, run, push, submit', 'Container lifecycle'],
        ['Deploy', 'deploy, add-ic, status, clean', 'SageMaker infrastructure'],
        ['Observe', 'test, logs, benchmark', 'Validation & monitoring'],
        ['Adapt', 'adapter (add/list/remove/update)', 'LoRA hot-swap'],
        ['Operate', 'register, export, validate, optimize', 'Day-2 operations'],
        ['Train', 'tune, train', 'Fine-tuning & bespoke training'],
        ['CI', 'ci, manifest', 'Pipeline integration'],
    ],
    col_widths=[2.0, 5.0, 4.7]
)
text_box(s, 0.8, 5.8, 11, 0.5, 'Self-contained bash with shared helpers in do/lib/. No Python runtime, no npm, no MCC dependency.', size=15, color=GRAY)
speaker_notes(s, 'A customer can rm -rf node_modules after generation. The project is fully standalone.\n\nThe do/ pattern is inspired by do-framework \u2014 simple, discoverable, composable.\n\nKey message: these scripts encode all the SageMaker API knowledge so the customer doesn\'t have to.')

# ===========================================================================
# SLIDE 8: Inference Components + Multi-IC
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Inference Components + Multi-IC', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'IC-Based Deployment as the Default', size=18, color=GRAY)
add_bullet_frame(s, 0.8, 1.8, 11, 1.5, [
    './do/deploy          \u2192 Creates endpoint + endpoint config + inference component',
    './do/add-ic          \u2192 Add a second model to the same endpoint',
    './do/status          \u2192 Check IC health and copy count',
], size=15)
text_box(s, 0.8, 3.5, 11, 0.4, 'What MCC handles automatically:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 4.0, 11, 2.5, [
    '\u2022  RoutingConfig with LEAST_OUTSTANDING_REQUESTS (required but undocumented)',
    '\u2022  INFERENCE_COMPONENT_NAME injection for log routing',
    '\u2022  GPU count allocation per IC',
    '\u2022  Configurable startup timeout (IC_STARTUP_TIMEOUT)',
    '\u2022  Container environment variable injection',
], size=15)
speaker_notes(s, 'Without RoutingConfig, ICs never schedule onto the endpoint. This is not documented anywhere \u2014 we discovered it through trial and error.\n\nMulti-IC is the future of SageMaker endpoints. One endpoint, multiple models, independent scaling.\n\nThe INFERENCE_COMPONENT_NAME injection is needed because SageMaker\'s IC log routing is broken \u2014 our CW forwarder uses it to write to the correct log group.')

# ===========================================================================
# SLIDE 9: LoRA Adapter Lifecycle
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'LoRA Adapter Lifecycle', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'Hot-Swap Adapters Without Redeploying', size=18, color=GRAY)
add_bullet_frame(s, 0.8, 1.8, 11, 2.0, [
    './do/adapter add --name=customer-finance --path=s3://bucket/adapters/finance-v2/',
    './do/adapter list',
    './do/adapter remove --name=customer-finance',
    './do/adapter update --name=customer-finance --path=s3://bucket/adapters/finance-v3/',
], size=14)
text_box(s, 0.8, 3.8, 11, 0.4, 'How it works:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 4.3, 11, 2.5, [
    '\u2022  vLLM --enable-lora + --max-loras=30 + --max-lora-rank=64',
    '\u2022  Adapters loaded at inference time via OpenAI-compatible API',
    '\u2022  do/adapters/ directory stores adapter configs for reproducibility',
    '\u2022  No redeploy, no downtime \u2014 adapter loaded on next request',
], size=15)
speaker_notes(s, 'A customer can train 30 LoRA adapters for different use cases and serve them all from one endpoint.\n\nNo redeploy, no downtime. The adapter is loaded on the next request that references it.\n\nUse case: one base model (Llama 3.1 8B) with adapters for finance, legal, code, customer support \u2014 all on one g5.2xlarge.')

# ===========================================================================
# SLIDE 10: Benchmarking
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Benchmarking', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'do/benchmark \u2014 Powered by SageMaker AI Benchmarking (NVIDIA AIPerf)', size=18, color=GRAY)
text_box(s, 0.8, 1.8, 11, 0.4, 'Example output:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 2.3, 6.0, 2.5, [
    'Request throughput:      12.4 req/s',
    'Output token throughput: 3,720 tok/s',
    'TTFT P50:               142ms',
    'TTFT P99:               891ms',
    'ITL P50:                28ms',
    'ITL P99:                67ms',
], size=14, color=GRAY)
text_box(s, 7.5, 1.8, 5.0, 0.4, 'Configurable:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 7.5, 2.3, 5.0, 2.5, [
    '\u2022  Concurrency (1\u2013100+)',
    '\u2022  Input/output token distribution',
    '\u2022  Streaming vs non-streaming',
    '\u2022  Request count',
    '\u2022  S3 output path for results',
], size=14)
text_box(s, 0.8, 5.2, 11, 0.4, 'Workflow:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 5.6, 11, 1.5, [
    'Generate with --include-benchmark \u2192 ./do/deploy \u2192 ./do/benchmark \u2192 ./do/benchmark --clean',
    'Compare instance types, model versions, and serving backends with real numbers.',
], size=14, color=GRAY)
speaker_notes(s, 'This is the same benchmarking service that powers SageMaker\'s internal performance testing.\n\nCustomers can compare instance types, model versions, and serving backends with real numbers.\n\nKey metrics: TTFT (time to first token) for user experience, ITL (inter-token latency) for streaming quality, throughput for cost efficiency.')

# ===========================================================================
# SLIDE 11: Validation & Registry
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Validation & Registry', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 5.5, 0.4, 'do/validate \u2014 Catch misconfigs before deploy:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 1.6, 6.0, 2.5, [
    '\u2705 Instance ml.g5.12xlarge has 4\u00d7 A10G (24GB each)',
    '\u2705 Model google/gemma-4-31B-it fits with TP=4',
    '\u2705 vLLM v0.20.2 supports Gemma4ForConditionalGeneration',
    '\u274c max_model_len=262144 exceeds available KV cache',
    '   \u2192 Recommendation: set VLLM_MAX_MODEL_LEN=16384',
], size=13)
text_box(s, 0.8, 4.3, 5.5, 0.4, 'do/register \u2014 Capture and replay:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 4.8, 6.0, 2.0, [
    './do/register                              # Capture success',
    'ml-container-creator registry replay prod  # Replay later',
    'ml-container-creator registry export       # Share with team',
], size=13, color=GRAY)
text_box(s, 7.5, 1.1, 5.0, 0.4, 'What gets validated:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 7.5, 1.6, 5.0, 3.0, [
    '\u2022  GPU memory vs model size',
    '\u2022  CUDA version compatibility',
    '\u2022  Instance availability in region',
    '\u2022  Model/server architecture match',
    '\u2022  Tensor parallelism feasibility',
    '\u2022  Context length vs available KV cache',
], size=13)
speaker_notes(s, 'Validation catches the max_model_len=262K OOM before you wait 15 minutes for a deploy to fail.\n\nThe registry is your team\'s institutional memory for deployments. Six months later, someone needs to redeploy \u2014 the registry has the exact config.\n\nExport as JSON, commit to git, import on any machine.')

# ===========================================================================
# SLIDE 12: CI Integration
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'CI Integration', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'From local dev to automated pipeline', size=18, color=GRAY)
text_box(s, 0.8, 1.8, 5.5, 0.4, 'Non-interactive generation:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 2.3, 6.0, 2.5, [
    'export MCC_DEPLOYMENT_CONFIG=transformers-vllm',
    'export MCC_MODEL_NAME=meta-llama/Llama-3.1-8B-Instruct',
    'export MCC_INSTANCE_TYPE=ml.g5.2xlarge',
    'ml-container-creator my-model --skip-prompts',
    '',
    'cd my-model',
    './do/submit && ./do/deploy && ./do/test && ./do/benchmark',
], size=13, color=GRAY)
text_box(s, 7.5, 1.8, 5.0, 0.4, 'CI Harness (bootstrap --ci):', size=16, color=DARK, bold=True)
add_bullet_frame(s, 7.5, 2.3, 5.0, 2.5, [
    '\u2022  CodeBuild project provisioned',
    '\u2022  E2E validation runner',
    '\u2022  Tests full lifecycle across configs',
    '\u2022  Runs on every PR',
    '\u2022  Catches regressions before customers',
], size=13)
text_box(s, 0.8, 5.5, 11, 0.4, 'Every parameter can be set via env var. No interactive prompts in CI.', size=15, color=ORANGE, bold=True)
speaker_notes(s, 'This slide matters most for platform teams and DevOps leads.\n\nThe CI harness is how we validate that vLLM v0.20.2 actually works before adding it to the catalog.\n\nScenario: ML platform team creates a config file with approved instance types. Data scientists run the generator with that config. In CI, use --skip-prompts with env vars.')

# ===========================================================================
# SLIDE 13: MCP Servers
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'MCP Servers \u2014 The Intelligence Layer', size=32, color=DARK, bold=True)
add_table(s, 0.8, 1.3, 11.7, 3.5,
    ['Server', 'What It Knows'],
    [
        ['instance-sizer', '"Gemma 4 31B needs 62GB VRAM \u2192 TP=4 on g5.12xlarge"'],
        ['region-picker', '"us-west-2 has g5.12xlarge quota, ap-southeast-1 doesn\'t"'],
        ['base-image-picker', '"vLLM v0.20.2 needs CUDA 12.9, use compat shim on driver 535"'],
        ['model-picker', '"google/gemma-4-31B-it is gated, needs HF_TOKEN"'],
        ['hyperpod-cluster-picker', '"Cluster ml-prod has 4 idle g5.12xlarge nodes"'],
        ['endpoint-picker', '"Endpoint shared-llm has capacity for another IC"'],
    ],
    col_widths=[3.5, 8.2]
)
text_box(s, 0.8, 5.3, 11, 0.4, 'Two modes:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 5.7, 6.0, 1.0, [
    'Static \u2014 Catalog data, no AWS credentials needed',
    'Smart \u2014 Live AWS API calls for real-time availability',
], size=14, color=GRAY)
text_box(s, 7.5, 5.3, 5.0, 0.4, 'AI Agent Integration:', size=16, color=DARK, bold=True)
add_bullet_frame(s, 7.5, 5.7, 5.0, 1.0, [
    'Kiro + MCP servers = autonomous deployment',
    '"Deploy Llama 3.1 to cheapest instance" \u2192 done',
], size=14, color=GRAY)
speaker_notes(s, 'MCP is the bridge between AI agents and deployment infrastructure.\n\nAt re:Invent 2026, imagine: "Hey Kiro, deploy this model" \u2192 MCC generates, builds, deploys, benchmarks.\n\nStatic mode requires zero AWS credentials \u2014 works from curated JSON catalogs shipped with the tool.')

# ===========================================================================
# SLIDE 14: LIVE DEMO - Deploy from Zero
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
# Orange accent box
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()
text_box(s, 0.8, 0.4, 11, 0.8, '\U0001f534 LIVE DEMO \u2014 Deploy from Zero', size=36, color=WHITE, bold=True)
text_box(s, 0.8, 1.5, 11, 0.4, 'GPT-OSS-20B on SageMaker \u2014 from nothing to benchmarked endpoint', size=20, color=DARK)
add_bullet_frame(s, 0.8, 2.3, 11, 4.5, [
    '1. Generate the project:',
    '   ml-container-creator gpt-oss-demo --deployment-config=transformers-vllm \\',
    '     --model-name=openai/gpt-oss-20b --instance-type=ml.g5.12xlarge --skip-prompts',
    '',
    '2. Show generated structure:  find . -maxdepth 2 | sort',
    '',
    '3. Build:  ./do/build     (~2 min, pulls vLLM base image)',
    '',
    '4. Push:   ./do/push      (~1 min, pushes to ECR)',
    '',
    '5. Deploy: ./do/deploy    (switch to pre-staged endpoint)',
    '',
    '6. Test + Logs + Adapter + Benchmark (next slide)',
], size=14)
speaker_notes(s, 'Run generation + build + push live (~3 min total).\n\nSwitch to pre-staged endpoint for deploy/test/benchmark (saves 15 min wait).\n\nSay: "Deploy takes about 15 minutes, so let me show you one I prepared earlier."\n\nKey message: "4 commands from nothing to a benchmarked, adapter-ready endpoint."')

# ===========================================================================
# SLIDE 15: DEMO - Test + Logs + Adapter
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()
text_box(s, 0.8, 0.4, 11, 0.8, '\U0001f534 DEMO \u2014 Test + Logs + Adapter + Benchmark', size=36, color=WHITE, bold=True)
text_box(s, 0.8, 1.5, 11, 0.4, 'On the pre-staged GPT-OSS-20B endpoint', size=20, color=DARK)
add_bullet_frame(s, 0.8, 2.2, 11, 4.8, [
    '1. Test inference:',
    '   ./do/test',
    '   \u2192 POST /v1/chat/completions \u2192 {"choices":[{"message":{"content":"Hello!..."}}]}',
    '',
    '2. Stream logs (CW forwarder in action):',
    '   ./do/logs',
    '   \u2192 Shows vLLM startup, model loading, inference requests',
    '',
    '3. Add a LoRA adapter:',
    '   ./do/adapter add --name=code-assist --path=s3://mlcc-adapters/code-assist-v1/',
    '   ./do/adapter list  \u2192 code-assist (active)',
    '',
    '4. Run benchmark:',
    '   ./do/benchmark',
    '   \u2192 Throughput, TTFT, ITL metrics',
], size=13)
speaker_notes(s, 'Show the actual inference response \u2014 proves the endpoint is live.\n\nShow logs flowing in real-time \u2014 this is the CW forwarder we built to work around broken IC log routing.\n\nShow adapter hot-swap working \u2014 no redeploy needed.\n\nIf benchmark takes too long, show a previous result screenshot.')

# ===========================================================================
# SLIDE 16: Platform Paper Cuts
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Platform Paper Cuts', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'Issues We\'ve Found \u2014 Opportunities to Fix', size=18, color=GRAY)
add_table(s, 0.8, 1.7, 11.7, 4.0,
    ['Issue', 'Impact', 'Current Workaround'],
    [
        ['IC log routing broken', 'Zero visibility into failures', 'Custom cw_log_forwarder.py'],
        ['RoutingConfig undocumented', 'ICs never schedule without it', 'MCC adds it automatically'],
        ['CUDA compat shim needed', 'Crash on driver mismatch', 'cuda_compat.sh patches LD_LIBRARY_PATH'],
        ['NCCL hangs on PCIe (g6/L4)', 'Multi-GPU models never start', 'NCCL_P2P_DISABLE=1 (investigating)'],
        ['vLLM env var deprecation', 'Breaking change in v0.20.2', 'EXCLUDE_VARS list in serve script'],
        ['Startup timeout too low', 'Large models fail health check', 'Configurable IC_STARTUP_TIMEOUT'],
    ],
    col_widths=[3.5, 3.7, 4.5]
)
text_box(s, 0.8, 6.2, 11, 0.5, 'Every one of these is a customer escalation waiting to happen.', size=15, color=ORANGE, bold=True)
speaker_notes(s, 'Every one of these is a customer escalation waiting to happen.\n\nMCC encodes the workarounds, but the platform should fix the root causes.\n\nIC log routing is the biggest one \u2014 customers have zero visibility when things go wrong. Log groups get created but zero log streams appear.\n\nAsk: "Who owns IC log routing? Can we get a ticket filed?"')

# ===========================================================================
# SLIDE 17: What We Built to Work Around Them
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Battle-Tested Fixes in Every Generated Project', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.3, 11, 0.4, 'cw_log_forwarder.py', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 0.8, 1.8, 11, 1.2, [
    '\u2022  Reads from stdin (piped from serve script via exec >)',
    '\u2022  Writes to CloudWatch Logs via boto3, 1-second flush interval',
    '\u2022  Log group: /aws/sagemaker/InferenceComponents/<IC_NAME>',
], size=14, color=GRAY)
text_box(s, 0.8, 3.2, 11, 0.4, 'cuda_compat.sh', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 0.8, 3.7, 11, 1.2, [
    '\u2022  Detects driver version vs required CUDA toolkit version',
    '\u2022  If mismatch: adds /usr/local/cuda/compat to LD_LIBRARY_PATH',
    '\u2022  Bridges driver 535 \u2192 CUDA 12.9 (well-tested NVIDIA path)',
], size=14, color=GRAY)
text_box(s, 0.8, 5.1, 11, 0.4, 'code/serve', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 0.8, 5.6, 11, 1.5, [
    '\u2022  Boolean flags: true \u2192 flag only, false \u2192 skip entirely',
    '\u2022  EXCLUDE_VARS: internal env vars that aren\'t CLI args',
    '\u2022  Model resolution: HuggingFace, S3, or local path',
], size=14, color=GRAY)
speaker_notes(s, 'These aren\'t hacks \u2014 they\'re production-grade workarounds that every customer needs.\n\nThe CW forwarder alone has saved us dozens of hours of blind debugging.\n\nThe boolean flag fix: vLLM --enable-lora takes no value. Passing "true" crashes it. Our serve script handles this correctly for all boolean VLLM_ env vars.')

# ===========================================================================
# SLIDE 18: re:Invent 2025 Recap
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 're:Invent 2025 Recap', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.5, 'Breakout session \u2014 generated + deployed a model live on stage', size=20, color=ORANGE, bold=True)
text_box(s, 0.8, 2.0, 11, 0.4, 'Since then (v0.3.0 \u2192 v0.5.0):', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 2.5, 5.5, 4.5, [
    '\u2705  4 deployment targets (was 1)',
    '\u2705  6 MCP servers (was 0)',
    '\u2705  Multi-IC endpoints',
    '\u2705  LoRA adapter lifecycle',
    '\u2705  SageMaker AI Benchmarking',
    '\u2705  Fine-tuning (managed + bespoke)',
    '\u2705  Marketplace model packages',
    '\u2705  CI integration harness',
], size=14)
add_bullet_frame(s, 7.0, 2.5, 5.5, 4.5, [
    '\u2705  Secrets Manager integration',
    '\u2705  Instance quota & availability',
    '\u2705  Schema-driven validation',
    '\u2705  Deployment registry',
    '\u2705  Post-deploy guidance',
    '\u2705  Yeoman removal (standalone CLI)',
    '\u2705  E2E validation runner',
    '\u2705  15 deployment configs (was 2)',
], size=14)
text_box(s, 0.8, 6.5, 11, 0.5, '16+ major features shipped. Weekly releases. 40+ specs designed and implemented.', size=15, color=GRAY)
speaker_notes(s, 'At re:Invent 2025, we had a CLI that generated a Dockerfile and a deploy script.\n\nToday we have a complete ML deployment platform in a box.\n\nThe velocity is real \u2014 weekly releases, 40+ specs designed and implemented.\n\nAudience reaction at re:Invent: "Why doesn\'t this exist already?"')

# ===========================================================================
# SLIDE 19: re:Invent 2026 Potential
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 're:Invent 2026 Potential', size=32, color=DARK, bold=True)
add_table(s, 0.8, 1.2, 11.7, 2.5,
    ['Format', 'Pitch'],
    [
        ['Workshop (2hr)', '"Deploy 3 models in 2 hours: LLM, diffusion, multi-model with LoRA"'],
        ['Chalk Talk', '"BYOC best practices: what we learned deploying 50 models"'],
        ['Builder Session', '"AI-powered deployment with MCC + Kiro + MCP"'],
        ['Breakout', '"From model weights to production inference in 5 minutes"'],
    ],
    col_widths=[3.0, 8.7]
)
text_box(s, 0.8, 4.2, 11, 0.4, 'Demo scenarios:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 4.7, 11, 2.5, [
    '\u2022  "Deploy Gemma 4 in 5 minutes" (latest model, day-0 support)',
    '\u2022  "Multi-model endpoint with LoRA hot-swap" (cost efficiency story)',
    '\u2022  "AI agent deploys a model" (Kiro + MCP \u2192 fully automated)',
    '\u2022  "Benchmark and optimize" (data-driven instance selection)',
], size=15)
text_box(s, 0.8, 6.5, 11, 0.5, 'The narrative: MCC is the missing layer between "I have model weights" and "I have a production endpoint."', size=15, color=ORANGE, bold=True)
speaker_notes(s, 'The workshop format is strongest \u2014 hands-on, attendees leave with a deployed model.\n\nThe AI agent demo is the wow factor \u2014 "Hey Kiro, deploy Llama 3.1 to the cheapest instance in us-west-2."\n\nWe need: a session slot, a demo account with pre-provisioned capacity, and a backup plan if live deploy fails.')

# ===========================================================================
# SLIDE 20: Roadmap
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'Roadmap', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.0, 11, 0.4, 'Shipped \u2705 \u2192 Designed \U0001f4d0 \u2192 Planned \U0001f4cb', size=18, color=GRAY)
add_table(s, 0.8, 1.6, 11.7, 4.5,
    ['Feature', 'Status', 'Description'],
    [
        ['Model architecture validation', '\U0001f4d0 Designed', 'Catch incompatible model/server combos at generation time'],
        ['E2E validation runner', '\U0001f4d0 In progress', 'Automated lifecycle testing across all 15 configs'],
        ['Notebook export', '\U0001f4d0 Designed', 'do/export --notebook for Jupyter sharing'],
        ['Cost estimation', '\U0001f4cb Planned', 'Estimated monthly cost before deploy'],
        ['Multi-region deployment', '\U0001f4cb Planned', 'Deploy same model to multiple regions'],
        ['Deployment manifest', '\U0001f4cb Planned', 'Single source of truth for all config values'],
        ['v1.0 release', '\U0001f4cb Planned', 'Stable API, npm publish, SLA on catalog updates'],
    ],
    col_widths=[3.5, 2.0, 6.2]
)
text_box(s, 0.8, 6.5, 11, 0.5, '30+ features shipped since re:Invent 2025. Weekly releases continue.', size=15, color=GRAY)
speaker_notes(s, 'We\'ve shipped 30+ features since re:Invent 2025. The tool is actively developed.\n\nModel architecture validation will catch "this model doesn\'t work on this server" at generation time.\n\nAsk: "Which of the planned features would be most valuable for your team?"')

# ===========================================================================
# SLIDE 21: The Vision
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'The Vision', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.2, 5.5, 0.4, 'Today:', size=20, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 1.7, 5.5, 2.5, [
    '\u2022  Customer has model weights',
    '\u2022  2-5 days of boilerplate + debugging',
    '\u2022  Maybe works',
    '\u2022  No benchmarks, no adapters, no CI',
    '\u2022  Knowledge trapped in one engineer',
], size=15, color=GRAY)
text_box(s, 7.0, 1.2, 5.5, 0.4, 'Future:', size=20, color=ORANGE, bold=True)
add_bullet_frame(s, 7.0, 1.7, 5.5, 2.5, [
    '\u2022  "Hey Kiro, deploy this model"',
    '\u2022  MCC + MCP generates + validates',
    '\u2022  Deployed, tested, benchmarked in 5 min',
    '\u2022  Adapters, CI, registry from day 1',
    '\u2022  Reproducible by anyone on the team',
], size=15)
text_box(s, 0.8, 4.5, 11, 0.4, 'Three pillars:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 5.0, 11, 2.0, [
    '1. Platform fixes \u2014 IC logging, RoutingConfig docs, NCCL \u2192 customers stop hitting these',
    '2. MCC maturity \u2014 v1.0, npm publish, community contributions, model registry',
    '3. AI integration \u2014 MCP servers + Kiro/Q = autonomous deployment',
], size=15)
speaker_notes(s, 'The end state is: deploying a model on SageMaker is as easy as deploying a Lambda function.\n\nMCC is the bridge from today\'s fragile manual process to tomorrow\'s AI-automated deployment.\n\nAll three pillars need to happen in parallel. Platform fixes reduce the workarounds MCC needs. MCC maturity increases adoption. AI integration is the multiplier.')

# ===========================================================================
# SLIDE 22: Call to Action
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)
text_box(s, 0.8, 0.4, 11, 0.6, 'What We Need From You', size=32, color=DARK, bold=True)
add_table(s, 0.8, 1.2, 11.7, 2.5,
    ['Ask', 'Why', 'How'],
    [
        ['Adopt', 'Validate with real customers', 'npx @aws/ml-container-creator \u2014 try it today'],
        ['Fix platform issues', 'Customers hitting these daily', 'IC logging, RoutingConfig docs, NCCL'],
        ['Fund', 'v1.0 needs headcount', 're:Invent session slot, team allocation'],
        ['Amplify', 'More users = more battle-testing', 'Share with teams, add to SA onboarding'],
    ],
    col_widths=[2.0, 4.0, 5.7]
)
text_box(s, 0.8, 4.2, 11, 0.4, 'Immediate asks:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 4.7, 11, 2.5, [
    '1.  IC log routing fix \u2014 who owns this? Can we get a ticket?',
    '2.  RoutingConfig documentation \u2014 it\'s required but undocumented',
    '3.  re:Invent 2026 session submission \u2014 workshop or builder session',
    '4.  Beta customers \u2014 who has a customer deploying BYOC this quarter?',
], size=15)
speaker_notes(s, 'Be specific: "I need a name for who owns IC log routing. I need a commitment to document RoutingConfig."\n\n"If you have a customer struggling with BYOC, send them to me. I\'ll generate their project in 5 minutes."\n\nLeave time for discussion on each ask.')

# ===========================================================================
# SLIDE 23: Thank You
# ===========================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
orange_bar(s, 7.2)
text_box(s, 1, 1.5, 11, 0.8, 'Try It Now', size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1, 2.8, 11, 0.8, 'npx @aws/ml-container-creator', size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_frame(s, 3, 4.0, 7, 2.0, [
    'GitHub:  github.com/awslabs/ml-container-creator',
    'Docs:    awslabs.github.io/ml-container-creator',
    'npm:     @aws/ml-container-creator',
], size=16, color=RGBColor(0xCC, 0xCC, 0xCC))
text_box(s, 1, 6.0, 11, 0.8, 'Questions?', size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
speaker_notes(s, 'Leave this slide up during Q&A.\n\nCommon questions: security review process, multi-region deployment, cost estimation, team onboarding.\n\nOffer: "I\'ll stay after to help anyone generate a project for their customer."\n\nReiterate: "5 minutes to try, no AWS credentials needed for local build+test."')

# ===========================================================================
# Save
# ===========================================================================
out = 'docs/deep-dive-deck.pptx'
prs.save(out)
print(f'\u2705 Generated {out} \u2014 {len(prs.slides)} slides')
