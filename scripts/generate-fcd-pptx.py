#!/usr/bin/env python3
"""POC: Generate an editable PPTX First-Call Deck for ML Container Creator."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# AWS brand-ish colors
DARK = RGBColor(0x23, 0x2F, 0x3E)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, left, top, width, height, text, size=18, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_frame(slide, left, top, width, height, items, size=16,
                     color=DARK, bold_first=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
        p.level = 0
    return tf

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'Calibri'

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
                p.font.name = 'Calibri'
    return tbl

def orange_bar(slide, top=0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top),
        prs.slide_width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ---------------------------------------------------------------------------
# SLIDE 1: Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s, DARK)
orange_bar(s, 7.2)

text_box(s, 1, 1.5, 11, 1.2, 'ML Container Creator', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1, 2.8, 11, 0.8, 'Deploy ML models to Amazon SageMaker — in minutes, not days', size=24, color=ORANGE, align=PP_ALIGN.CENTER)
text_box(s, 1, 4.2, 11, 0.5, 'npx @aws/generator-ml-container-creator', size=18, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)
text_box(s, 1, 5.2, 11, 0.5, 'Open Source  ·  Apache 2.0  ·  github.com/awslabs/ml-container-creator', size=14, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

speaker_notes(s, 'Introduce yourself and the tool in one sentence: "ML Container Creator is an open-source code generator that takes your ML model and produces a complete, deployable SageMaker container project in under a minute." Mention it\'s an AWS Labs project — community-driven, Apache 2.0 licensed, not a managed service.')

# ---------------------------------------------------------------------------
# SLIDE 2: The Problem
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'The Problem', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.2, 11, 0.6, 'Getting a model from a notebook to a SageMaker endpoint is harder than it should be.', size=20, color=GRAY)

add_bullet_frame(s, 0.8, 2.2, 11, 3.5, [
    '1.  A Dockerfile meeting SageMaker BYOC requirements (port 8080, /ping, /invocations)',
    '2.  A model server choice and configuration (Flask? vLLM? Triton?)',
    '3.  Model loading and inference code',
    '4.  Deployment scripts (ECR push, endpoint creation, IAM roles)',
    '5.  Testing infrastructure (local container tests, endpoint tests)',
], size=16)

text_box(s, 0.8, 5.5, 11, 0.5, 'The result:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 6.0, 11, 1.2, [
    '•  Teams repeat this boilerplate for every model × framework × server',
    '•  Container misconfigurations → failed deployments, wasted GPU hours',
    '•  No standardization → knowledge silos across teams',
], size=14, color=GRAY)

speaker_notes(s, 'Ask the audience: "How many of you have a Dockerfile for SageMaker that you copy-paste between projects?" — this is the pain point.\n\nEmphasize that the problem isn\'t any single step — it\'s the combination of all of them, and the fact that each framework has completely different container requirements.')

# ---------------------------------------------------------------------------
# SLIDE 3: What MCC Does
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'What ML Container Creator Does', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.3, 11, 0.6, 'One CLI command → complete, buildable project', size=22, color=ORANGE, bold=True)

add_bullet_frame(s, 0.8, 2.3, 11, 3.0, [
    '✅  SageMaker-compatible Dockerfile (framework-specific)',
    '✅  Model serving code (handler, server, configs)',
    '✅  Lifecycle scripts (build, push, deploy, test, clean, logs)',
    '✅  Optional sample model for immediate testing',
    '✅  Optional test suite + IAM permission reference',
], size=16)

text_box(s, 0.8, 5.2, 11, 0.5, 'Key principle', size=18, color=DARK, bold=True)
text_box(s, 0.8, 5.7, 11, 1.0, "It's a code generator, not a runtime framework.\nYou own the output. No agent in your container. No lock-in.", size=16, color=GRAY)

speaker_notes(s, 'Stress the "code generator" distinction. This is not like SageMaker Inference Toolkit or a framework you import. It generates plain files — Dockerfiles, Python scripts, bash scripts — that you can read, modify, and commit to your repo.\n\nThe generated code is starter code. We don\'t claim "production-ready" — we say "SageMaker-compatible."')

# ---------------------------------------------------------------------------
# SLIDE 4: Supported Architectures
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'Supported Architectures', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, '4 architecture families · 15 deployment configurations', size=18, color=GRAY)

add_table(s, 0.8, 1.8, 11.7, 4.5,
    ['Architecture', 'Deployment Configs', 'Use Case'],
    [
        ['HTTP (2)', 'http-flask, http-fastapi', 'Traditional ML (sklearn, XGBoost, TF) · CPU · ms latency'],
        ['Transformers (5)', 'transformers-vllm, -sglang,\n-tensorrt-llm, -lmi, -djl', 'LLM serving · HuggingFace · GPU · seconds latency'],
        ['Triton (7)', 'triton-fil, -onnxruntime, -tensorflow,\n-pytorch, -vllm, -tensorrtllm, -python', 'NVIDIA Triton · high-throughput · multi-model'],
        ['Diffusors (1)', 'diffusors-vllm-omni', 'Image generation (SD, FLUX) · roadmap'],
    ],
    col_widths=[2.2, 4.5, 5.0]
)

speaker_notes(s, 'Walk through each architecture family briefly. The key insight is that these aren\'t just different frameworks — they have fundamentally different container architectures:\n- HTTP: your code loads the model, Nginx sits in front\n- Transformers: the framework IS the server\n- Triton: NVIDIA\'s inference server with model repository pattern\n\nAsk: "Which of these architectures are you currently using or evaluating?"')

# ---------------------------------------------------------------------------
# SLIDE 5: MCP Servers
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'Intelligent Defaults via MCP Servers', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, '5 bundled advisors that help you make better choices', size=18, color=GRAY)

add_table(s, 0.8, 1.8, 11.7, 3.0,
    ['Server', 'What It Does'],
    [
        ['🖥️  Instance Recommender', 'Right-sized SageMaker instances for your framework and model'],
        ['🌍  Region Picker', 'Filter AWS regions by availability and proximity'],
        ['📦  Base Image Picker', 'Curated, versioned container images per framework'],
        ['☸️  HyperPod Cluster Picker', 'Discover existing HyperPod EKS clusters via AWS API'],
        ['🤖  Model Picker', 'Resolve HuggingFace model metadata (architecture, gated status)'],
    ],
    col_widths=[3.5, 8.2]
)

add_bullet_frame(s, 0.8, 5.3, 11, 1.5, [
    'Static mode (default) — instant responses from curated catalogs, no AWS credentials needed',
    'Smart mode (opt-in) — Amazon Bedrock-powered context-aware recommendations',
    'MCP servers are configuration providers, not AI agents. No LLM in the loop unless you opt in.',
], size=14, color=GRAY)

speaker_notes(s, 'Explain MCP simply: "These are small helper programs that the generator talks to over stdio. They answer questions like \'what instance types work well for vLLM?\' and the generator uses those answers to populate your choices."\n\nStatic mode requires zero AWS credentials — it works from curated JSON catalogs shipped with the tool.')

# ---------------------------------------------------------------------------
# SLIDE 6: CI/CD Configuration
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'Configuration for CI/CD', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, '8-level precedence — prompts are the last resort', size=18, color=GRAY)

add_bullet_frame(s, 0.8, 1.8, 5.5, 2.5, [
    'CLI Options (highest priority)',
    '  → CLI Arguments',
    '    → Environment Variables',
    '      → CLI Config File (--config)',
    '        → Custom Config File',
    '          → Package.json section',
    '            → Generator Defaults',
    '              → Interactive Prompts (lowest)',
], size=13)

# Right column: example
text_box(s, 7, 1.8, 5.5, 0.4, 'Fully automated:', size=14, color=DARK, bold=True)
add_bullet_frame(s, 7, 2.3, 5.5, 2.0, [
    'yo @aws/ml-container-creator \\',
    '  --skip-prompts \\',
    '  --deployment-config=transformers-vllm \\',
    '  --model-name=meta-llama/Llama-2-7b-chat-hf \\',
    '  --instance-type=ml.g5.xlarge',
], size=12, color=GRAY)

text_box(s, 7, 4.5, 5.5, 0.4, 'Environment variables:', size=14, color=DARK, bold=True)
add_bullet_frame(s, 7, 4.9, 5.5, 1.5, [
    'ML_INSTANCE_TYPE=ml.g5.xlarge',
    'AWS_REGION=us-east-1',
    'AWS_ROLE=arn:aws:iam::...:role/SageMakerRole',
], size=12, color=GRAY)

speaker_notes(s, 'This slide matters most for platform teams and DevOps leads. The interactive prompts are great for exploration, but real adoption happens when you can run this in a CI pipeline.\n\nScenario: "Your ML platform team creates a config file with approved instance types and regions. Individual data scientists run the generator with that config file. In CI, you use --skip-prompts with environment variables."')

# ---------------------------------------------------------------------------
# SLIDE 7: Deployment Targets
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'Deployment Targets', size=32, color=DARK, bold=True)
text_box(s, 0.8, 1.1, 11, 0.4, 'Same container, multiple deployment paths', size=18, color=GRAY)

# Left: Managed Inference
text_box(s, 0.8, 1.8, 5.5, 0.4, 'SageMaker Managed Inference', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 0.8, 2.4, 5.5, 2.5, [
    '•  Real-time endpoints via Inference Components API',
    '•  ./do/deploy <role-arn>',
    '•  ./do/test <endpoint-name>',
    '•  ./do/logs → CloudWatch tailing',
    '•  ./do/clean endpoint → tear down',
], size=14)

# Right: HyperPod
text_box(s, 7, 1.8, 5.5, 0.4, 'SageMaker HyperPod EKS', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 7, 2.4, 5.5, 2.5, [
    '•  K8s manifests: Deployment, Service, ConfigMap, PVC',
    '•  ./do/deploy → kubectl apply',
    '•  FSx for Lustre volume support',
    '•  Cluster discovery via MCP server',
    '•  Same container image as managed inference',
], size=14)

# Bottom: Build targets
text_box(s, 0.8, 5.2, 11, 0.4, 'Build Targets', size=18, color=ORANGE, bold=True)
add_bullet_frame(s, 0.8, 5.7, 11, 1.0, [
    'Local — ./do/build + ./do/push (Docker on your machine)',
    'CodeBuild — ./do/submit (cloud-based build, no local Docker needed — critical for 10-20GB LLM containers)',
], size=14)

speaker_notes(s, 'The key message: the container image is the same regardless of deployment target. What changes is the deployment script and the infrastructure manifests.\n\nCodeBuild is important for LLM containers that can be 10-20GB — building those locally is painful.')

# ---------------------------------------------------------------------------
# SLIDE 8: What It's Not
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, "What It's Not", size=32, color=DARK, bold=True)

add_table(s, 0.8, 1.3, 11.7, 2.5,
    ['✅ It Is', '❌ It Is Not'],
    [
        ['A code generator', 'A managed service'],
        ['Starter code you own and modify', 'Production-ready without review'],
        ['SageMaker BYOC tooling', 'A replacement for JumpStart'],
        ['Framework-agnostic scaffolding', 'A runtime framework (no lock-in)'],
        ['Open source (Apache 2.0)', 'An AWS service with SLA'],
    ],
    col_widths=[5.85, 5.85]
)

text_box(s, 0.8, 4.5, 11, 0.4, 'When to use something else:', size=18, color=DARK, bold=True)
add_bullet_frame(s, 0.8, 5.0, 11, 2.0, [
    '•  Model supported by built-in algorithms → use those',
    '•  Want one-click deployment of a popular model → use JumpStart',
    '•  Need managed MLOps pipeline → use SageMaker Pipelines',
    '•  Need full container control, custom serving, unsupported frameworks → use MCC ✅',
], size=15)

speaker_notes(s, 'This slide builds trust. Being honest about limitations makes the strengths more credible.\n\nMost common objection: "why not just use JumpStart?" — the answer is control. JumpStart is great when it supports your model. MCC is for when you need to customize the container.')

# ---------------------------------------------------------------------------
# SLIDE 9: Roadmap
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
orange_bar(s)

text_box(s, 0.8, 0.4, 11, 0.6, 'Roadmap', size=32, color=DARK, bold=True)

add_table(s, 0.8, 1.3, 11.7, 4.0,
    ['Feature', 'Status', 'Description'],
    [
        ['Async Inference Endpoints', 'Designed', 'Large payloads (1GB), long inference (1hr), scale-to-zero'],
        ['Diffusion Model Support', 'Designed', 'Image generation via vLLM-Omni (Stable Diffusion, FLUX)'],
        ['Triton Sample Models', 'Designed', 'Auto-training for Triton backends (FIL, ONNX, TF, Python)'],
        ['S3 Model Loading', 'Planned', 'Load models from S3 at runtime instead of baking into container'],
        ['JumpStart Integration', 'Planned', 'Use JumpStart model artifacts with custom containers'],
        ['Model Registry', 'Planned', 'Pull models from SageMaker Model Registry'],
    ],
    col_widths=[3.5, 1.5, 6.7]
)

text_box(s, 0.8, 5.8, 11, 0.8, 'Designed = full spec exists  ·  Planned = on roadmap, not yet specced\nContributions and use cases welcome!', size=14, color=GRAY)

speaker_notes(s, 'Gauge interest in roadmap items. If the audience is heavy on LLMs, emphasize async inference. If traditional ML, emphasize S3 model loading.\n\nInvite contributions: "If any of these are critical for your team, we\'d love contributions or even just detailed use cases to help prioritize."')

# ---------------------------------------------------------------------------
# SLIDE 10: Getting Started + Thank You
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
orange_bar(s, 7.2)

text_box(s, 1, 1.0, 11, 0.8, 'Getting Started', size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_bullet_frame(s, 2.5, 2.2, 8, 2.5, [
    'Prerequisites: Node.js 24+  ·  Python 3.8+  ·  Docker 20+  ·  AWS CLI 2+',
    '',
    'git clone https://github.com/awslabs/ml-container-creator.git',
    'cd ml-container-creator && npm install && npm link',
    'yo @aws/ml-container-creator',
], size=16, color=RGBColor(0xCC,0xCC,0xCC))

text_box(s, 1, 5.0, 11, 0.8, 'Thank You — Questions?', size=32, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
text_box(s, 1, 6.0, 11, 0.5, '📖 awslabs.github.io/ml-container-creator   ·   💻 github.com/awslabs/ml-container-creator', size=14, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

speaker_notes(s, 'Clear call to action: "Try generating a container for one of your existing models this week. If you hit issues, open a GitHub issue."\n\nThe Node.js 24+ requirement may surprise people — mention nvm for version management.')

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = 'docs/first-call-deck.pptx'
prs.save(out)
print(f'✅ Generated {out} — {len(prs.slides)} slides')
